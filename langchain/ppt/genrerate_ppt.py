import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

# Fix 1: Safe local output path
out = "LangChain_Architecture_Deep_Dive.pptx"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BLUE = RGBColor(47, 84, 150)
LIGHT = RGBColor(221, 235, 247)
GRAY = RGBColor(242, 242, 242)
DARK = RGBColor(55, 55, 55)
GREEN = RGBColor(226, 239, 218)
ORANGE = RGBColor(252, 228, 214)
PURPLE = RGBColor(234, 209, 220)

def title(s, t, sub):
    tb = s.shapes.add_textbox(Inches(.4), Inches(.18), Inches(12.5), Inches(.65))
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = t
    r.font.size = Pt(25)
    r.font.bold = True
    r.font.color.rgb = BLUE

    sb = s.shapes.add_textbox(Inches(.65), Inches(.76), Inches(12), Inches(.38))
    p = sb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = sub
    r.font.size = Pt(12)
    r.font.color.rgb = DARK

def box(s, x, y, w, h, text, fill=LIGHT, fs=15):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = BLUE
    sh.line.width = Pt(1.4)
    
    tf = sh.text_frame
    tf.word_wrap = True
    # Fix 2: Set margins to avoid clipping
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    lines = text.split('\n')
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = line
        r.font.size = Pt(fs)
        r.font.bold = (i == 0) # Bold heading line
        r.font.color.rgb = DARK

def arrow(s, x1, y1, x2, y2):
    l = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    l.line.color.rgb = BLUE
    l.line.width = Pt(2)
    l.line.end_arrowhead = True

def note(s, t):
    tb = s.shapes.add_textbox(Inches(.55), Inches(6.86), Inches(12.2), Inches(.3))
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = t
    r.font.size = Pt(11)
    r.font.italic = True
    r.font.color.rgb = DARK

# Slide 1
s = prs.slides.add_slide(prs.slide_layouts[6])
title(s, "LangChain Architecture", "From application input to model reasoning, tools, retrieval, state, and final output")
box(s, .4, 2.85, 1.65, 1, "User", GRAY, 16)
arrow(s, 2.05, 3.35, 2.55, 3.35)
box(s, 2.55, 2.25, 2.1, 2.1, "Application\nPrompt • State\nBusiness Logic", LIGHT, 15)
arrow(s, 4.65, 3.3, 5.15, 3.3)
box(s, 5.15, 2.05, 2.25, 2.5, "LangChain\nRunnables / LCEL\nChains / Agents", LIGHT, 15)
arrow(s, 7.4, 3.3, 7.9, 3.3)
box(s, 7.9, 2.25, 2.05, 2.1, "Model\nChat LLM\nEmbeddings", ORANGE, 15)
arrow(s, 9.95, 3.35, 10.5, 3.35)
box(s, 10.5, 2.85, 2.2, 1, "Response", GRAY, 16)
box(s, 3.1, 4.95, 2.0, 1, "Tools\nAPIs / DB", GREEN, 14)
arrow(s, 5.1, 5.45, 6.15, 4.5)
box(s, 6.15, 4.95, 2.0, 1, "RAG\nRetriever", GREEN, 14)
arrow(s, 7.15, 4.95, 7.15, 4.5)
box(s, 9.0, 4.95, 2.0, 1, "Memory /\nState", PURPLE, 14)
arrow(s, 9.95, 4.95, 9.25, 4.35)
note(s, "LangChain is an orchestration layer around models, data, state, tools, and application logic.")

# Slide 2
s = prs.slides.add_slide(prs.slide_layouts[6])
title(s, "Runnables: The Core Execution Abstraction", "Composable units that can be invoked, streamed, batched, and connected")
box(s, .5, 2.55, 1.7, 1.1, "Input", GRAY, 16)
arrow(s, 2.2, 3.1, 2.7, 3.1)
box(s, 2.7, 2.2, 2.0, 1.8, "Runnable\ninvoke()\nbatch()\nstream()", LIGHT, 15)
arrow(s, 4.7, 3.1, 5.2, 3.1)
box(s, 5.2, 2.2, 2.0, 1.8, "Prompt\nRunnable", LIGHT, 15)
arrow(s, 7.2, 3.1, 7.7, 3.1)
box(s, 7.7, 2.2, 2.0, 1.8, "Chat Model\nRunnable", ORANGE, 15)
arrow(s, 9.7, 3.1, 10.2, 3.1)
box(s, 10.2, 2.55, 2.45, 1.1, "Output", GRAY, 16)
box(s, 2.5, 5.0, 8.4, 1, "Composition: Prompt | Model | Parser | Application", GREEN, 16)
note(s, "A Runnable gives LangChain components a common execution model and makes composition predictable.")

# Slide 3
s = prs.slides.add_slide(prs.slide_layouts[6])
title(s, "LCEL — LangChain Expression Language", "Declaratively compose Runnables into readable pipelines")
box(s, .45, 2.55, 2.05, 1.55, "PromptTemplate\n{name} + {question}", LIGHT, 15)
arrow(s, 2.5, 3.32, 3.05, 3.32)
box(s, 3.05, 2.55, 2.05, 1.55, "Chat Model\nLLM call", ORANGE, 15)
arrow(s, 5.1, 3.32, 5.65, 3.32)
box(s, 5.65, 2.55, 2.05, 1.55, "Output Parser\nString / JSON", GREEN, 15)
arrow(s, 7.7, 3.32, 8.25, 3.32)
box(s, 8.25, 2.55, 2.2, 1.55, "Application\ninvoke / stream", GRAY, 15)
box(s, 2.15, 5, 8.9, 1, "Pipeline = Prompt → Model → Parser → Application", LIGHT, 16)
note(s, "LCEL makes execution graphs explicit and supports composition patterns such as sequential and parallel work.")

# Slide 4
s = prs.slides.add_slide(prs.slide_layouts[6])
title(s, "Chains: Predictable Multi-Step Workflows", "Connect deterministic steps when the workflow is known in advance")
box(s, .45, 2.65, 1.65, 1.15, "Input", GRAY, 16)
arrow(s, 2.1, 3.22, 2.55, 3.22)
box(s, 2.55, 2.25, 2.0, 1.95, "Step 1\nClassify /\nExtract", LIGHT, 15)
arrow(s, 4.55, 3.22, 5, 3.22)
box(s, 5, 2.25, 2, 1.95, "Step 2\nPrompt +\nModel", ORANGE, 15)
arrow(s, 7, 3.22, 7.45, 3.22)
box(s, 7.45, 2.25, 2, 1.95, "Step 3\nTransform /\nParse", GREEN, 15)
arrow(s, 9.45, 3.22, 9.9, 3.22)
box(s, 9.9, 2.65, 2.45, 1.15, "Output", GRAY, 16)
box(s, 3, 5.05, 7.5, 1, "Known workflow: Step 1 → Step 2 → Step 3 → Result", LIGHT, 15)
note(s, "Modern LangChain applications commonly express chains through Runnables and LCEL.")

# Slide 5
s = prs.slides.add_slide(prs.slide_layouts[6])
title(s, "Agents & Tools", "An agent lets the model decide whether and how to use available tools")
box(s, .35, 2.85, 1.55, 1.05, "User", GRAY, 16)
arrow(s, 1.9, 3.38, 2.35, 3.38)
box(s, 2.35, 2.3, 2.2, 2.15, "Agent\nInstructions\nModel\nTool Definitions", LIGHT, 15)
arrow(s, 4.55, 3.38, 5.05, 3.38)
box(s, 5.05, 2.35, 2.2, 2.05, "LLM\nAnswer or\nTool Call?", ORANGE, 15)
arrow(s, 7.25, 3, 8, 2)
box(s, 8, 1.25, 2.25, 1.2, "Tool 1\nSearch", GREEN, 14)
arrow(s, 7.25, 3.85, 8, 4.15)
box(s, 8, 3.55, 2.25, 1.2, "Tool 2\nDB / API", GREEN, 14)
arrow(s, 10.25, 1.85, 10.9, 3)
arrow(s, 10.25, 4.15, 10.9, 3.55)
box(s, 10.9, 3, 1.95, 1.2, "Observation", GRAY, 15)
box(s, 6.1, 5, 3.1, 1, "Repeat until\nfinal answer", LIGHT, 15)
note(s, "Agentic execution is iterative: decide → act → observe → decide again.")

# Slide 6
s = prs.slides.add_slide(prs.slide_layouts[6])
title(s, "Memory & State", "Conversation context and application state can be carried across model calls")
box(s, .55, 2.7, 1.8, 1.1, "User\nMessage", GRAY, 16)
arrow(s, 2.35, 3.25, 2.8, 3.25)
box(s, 2.8, 2.25, 2.15, 2, "Application State\nMessages /\nSession Data", PURPLE, 15)
arrow(s, 4.95, 3.25, 5.45, 3.25)
box(s, 5.45, 2.25, 2.2, 2, "Prompt Assembly\nHistory +\nNew Input", LIGHT, 15)
arrow(s, 7.65, 3.25, 8.15, 3.25)
box(s, 8.15, 2.25, 2.2, 2, "Model\nUses\nContext", ORANGE, 15)
arrow(s, 10.35, 3.25, 10.85, 3.25)
box(s, 10.85, 2.7, 1.9, 1.1, "Answer", GRAY, 16)
box(s, 3.1, 5.05, 7.4, 1, "Long-term knowledge is usually better handled with retrieval than unlimited chat history.", GREEN, 14)
note(s, "In modern architectures, state can be explicit and persisted; LangGraph adds durable, stateful execution for complex workflows.")

# Slide 7
s = prs.slides.add_slide(prs.slide_layouts[6])
title(s, "RAG — Retrieval-Augmented Generation", "Retrieve relevant external knowledge and place it into the model context")
box(s, .35, 2.65, 1.75, 1.1, "User Query", GRAY, 16)
arrow(s, 2.1, 3.2, 2.55, 3.2)
box(s, 2.55, 2.25, 2, 2, "Query Embedding\nText → Vector", LIGHT, 15)
arrow(s, 4.55, 3.2, 5, 3.2)
box(s, 5, 2.25, 2.2, 2, "Retriever\nSimilarity /\nMetadata Filter", GREEN, 15)
arrow(s, 7.2, 3.2, 7.65, 3.2)
box(s, 7.65, 2.25, 2.15, 2, "Relevant Documents\nTop-K Chunks", GREEN, 15)
arrow(s, 9.8, 3.2, 10.25, 3.2)
box(s, 10.25, 2.65, 2.5, 1.1, "Prompt + LLM", ORANGE, 16)
box(s, 4, 5.05, 5.6, 1, "Query + Retrieved Context → Grounded Generation", LIGHT, 16)
note(s, "RAG keeps external knowledge outside the model and retrieves relevant context at query time.")

# Slide 8
s = prs.slides.add_slide(prs.slide_layouts[6])
title(s, "Embeddings & Vector Stores", "Represent content as vectors so semantically similar information can be found")
box(s, .45, 2.45, 1.8, 1.25, "Documents", GRAY, 16)
arrow(s, 2.25, 3.08, 2.7, 3.08)
box(s, 2.7, 2.05, 2, 2.05, "Chunking\nSplit into\npassages", LIGHT, 15)
arrow(s, 4.7, 3.08, 5.15, 3.08)
box(s, 5.15, 2.05, 2, 2.05, "Embedding Model\nText → Vector", ORANGE, 15)
arrow(s, 7.15, 3.08, 7.6, 3.08)
box(s, 7.6, 2.05, 2.2, 2.05, "Vector Store\nVectors +\nMetadata", GREEN, 15)
arrow(s, 9.8, 3.08, 10.25, 3.08)
box(s, 10.25, 2.45, 2.55, 1.25, "Similarity Search", GRAY, 16)
box(s, 2.2, 5.05, 8.7, 1, "Query → Embedding → Search → Chunks → LLM", LIGHT, 15)
note(s, "Vector stores implement the retrieval layer; embeddings provide the numerical representation used for semantic search.")

# Slide 9
s = prs.slides.add_slide(prs.slide_layouts[6])
title(s, "LangGraph: Stateful Agent Workflows", "Explicit graphs with nodes, edges, state, branching, loops, and persistence")
box(s, .4, 2.75, 1.7, 1.05, "START", GRAY, 16)
arrow(s, 2.1, 3.28, 2.55, 3.28)
box(s, 2.55, 2.2, 2, 2.1, "Node: Agent\nRead State\nCall Model", LIGHT, 15)
arrow(s, 4.55, 3.28, 5, 3.28)
box(s, 5, 2.2, 2, 2.1, "Node: Tools\nExecute API /\nDatabase", GREEN, 15)
arrow(s, 7, 3.28, 7.45, 3.28)
box(s, 7.45, 2.2, 2, 2.1, "Node: Observe\nUpdate State", PURPLE, 15)
arrow(s, 9.45, 3.28, 10, 3.28)
box(s, 10, 2.75, 1.8, 1.05, "END", GRAY, 16)
box(s, 5.05, 5, 2.5, 1, "Conditional Edge\nContinue / Finish", LIGHT, 14)
box(s, 1, 5.75, 2.4, .8, "Shared State", PURPLE, 14)
note(s, "LangGraph is useful for durable state, cycles, branching, human-in-the-loop, and long-running agent workflows.")

# Slide 10
s = prs.slides.add_slide(prs.slide_layouts[6])
title(s, "Putting It All Together — LangChain + LangGraph", "A practical agentic RAG architecture")
box(s, .35, 2.85, 1.55, 1, "User", GRAY, 15)
arrow(s, 1.9, 3.35, 2.3, 3.35)
box(s, 2.3, 2.15, 1.95, 2.4, "LangGraph\nState + Flow\nAgent Nodes", LIGHT, 15)
arrow(s, 4.25, 3.35, 4.7, 3.35)
box(s, 4.7, 2.15, 1.85, 2.4, "LangChain\nRunnables\nLCEL / Chains", LIGHT, 15)
arrow(s, 6.55, 3.35, 7, 3.35)
box(s, 7, 2.15, 1.8, 2.4, "Model\nChat LLM\nEmbeddings", ORANGE, 15)
box(s, 9.15, 1.45, 1.8, 1.15, "RAG\nRetriever", GREEN, 14)
arrow(s, 8.8, 2.95, 9.15, 2.05)
box(s, 9.15, 3, 1.8, 1.15, "Tools\nAPIs / DB", GREEN, 14)
arrow(s, 8.8, 3.45, 9.15, 3.55)
box(s, 9.15, 4.55, 1.8, 1.15, "State /\nMemory", PURPLE, 14)
arrow(s, 8.8, 3.95, 9.15, 5.05)
box(s, 2, 5.85, 9, .75, "Request → State → Agent → Model → RAG / Tools → Observation → State → Answer", LIGHT, 14)
note(s, "Mental model: LangChain provides composable building blocks; LangGraph provides explicit stateful control flow for sophisticated agents.")

prs.save(out)
print(f"Saved successfully to {out}")